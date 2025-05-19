import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
import io

@st.cache_data(show_spinner=False)
def extract_text_from_additional(uploaded_file):
    """Extracts text from PDF or PPTX file."""
    if uploaded_file is None:
        return ""
    
    file_bytes = uploaded_file.getvalue()
    file_name = uploaded_file.name.lower()

    text = ""
    try:
        if file_name.endswith(".pdf"):
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text() + "\n"
        elif file_name.endswith(".pptx"):
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            st.warning(f"Unsupported file type for additional context: {uploaded_file.name}. Please use PDF or PPTX.")
            return ""
    except Exception as e:
        st.error(f"Error extracting text from {uploaded_file.name}: {e}")
        return ""
    return text